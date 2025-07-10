import os
from tkinter import Tk, Label, StringVar, Scrollbar, Text, Button, Frame
from tkinterdnd3 import DND_FILES, TkinterDnD
from pptx import Presentation
from logging_config import logger
import shutil
import zipfile
import tempfile
import xml.etree.ElementTree as ET
import zlib
from PIL import Image
import io
from tkinter import ttk
from PIL import ImageTk
import tkinter.messagebox as messagebox
from tkinter import Toplevel
from PIL import ImageDraw

def get_image_filename_from_zip(pptx_path, image_blob):
    """Trouve le nom du fichier image dans le ZIP en comparant les CRC."""
    try:
        # Calcul du CRC32 du blob de l'image
        image_crc = zlib.crc32(image_blob) & 0xFFFFFFFF
        
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            # Parcours des fichiers dans ppt/media
            for file_info in zip_ref.infolist():
                if file_info.filename.startswith('ppt/media/'):
                    # Si le CRC correspond, on a trouvé le fichier
                    if file_info.CRC == image_crc:
                        return os.path.basename(file_info.filename)
        
        return "Image sans nom"
    except Exception as e:
        logger.error(f"Erreur lors de la recherche du nom de fichier : {str(e)}")
        return "Image sans nom"

def get_images_from_shapes(shapes, pptx_path, slide_index=None):
    """Récupère le blob d'image et vérifie si l'image est rognée."""
    logger.debug(f"Analyse de {len(shapes)} formes")
    images = {}
    for i, shape in enumerate(shapes):
        logger.debug(f"Traitement de la forme {i+1}/{len(shapes)}")
        if hasattr(shape, "image"):
            image_blob = shape.image.blob
            if image_blob not in images:
                # Utilisation de la nouvelle fonction pour obtenir le nom du fichier
                filename = get_image_filename_from_zip(pptx_path, image_blob)
                logger.debug(f"Nouvelle image détectée, taille : {len(image_blob)/1024:.2f} KB, nom : {filename}")
                images[image_blob] = {
                    "used": False, 
                    "cropped": False,
                    "crop_ratio": 0.0,
                    "filename": filename,
                    "crop_left": 0.0,
                    "crop_top": 0.0,
                    "crop_right": 0.0,
                    "crop_bottom": 0.0,
                    "slide_index": slide_index
                }
            # Calcul du ratio de rognage
            if shape.crop_left > 0 or shape.crop_top > 0 or shape.crop_right > 0 or shape.crop_bottom > 0:
                logger.debug(f"Image rognée détectée - Crop values: L={shape.crop_left}, T={shape.crop_top}, R={shape.crop_right}, B={shape.crop_bottom}")
                images[image_blob]["cropped"] = True
                images[image_blob]["crop_left"] = shape.crop_left
                images[image_blob]["crop_top"] = shape.crop_top
                images[image_blob]["crop_right"] = shape.crop_right
                images[image_blob]["crop_bottom"] = shape.crop_bottom
                crop_ratio = (shape.crop_left + shape.crop_right + shape.crop_top + shape.crop_bottom) / 4
                images[image_blob]["crop_ratio"] = min(crop_ratio, 1.0)
                logger.debug(f"Image rognée - Nom: {images[image_blob]['filename']}, "
                           f"Taille: {len(image_blob)/1024:.2f} KB, "
                           f"Ratio de rognage: {images[image_blob]['crop_ratio']*100:.1f}%, "
                           f"Slide: {slide_index if slide_index is not None else 'N/A'}")
    return images

def extract_image(blob, filename, output_dir):
    """Extrait une image dans le répertoire spécifié."""
    logger.debug(f"Tentative d'extraction de l'image {filename} vers {output_dir}")
    try:
        if not os.path.exists(output_dir):
            logger.debug(f"Création du répertoire {output_dir}")
            os.makedirs(output_dir)
            
        output_path = os.path.join(output_dir, filename)
        logger.debug(f"Écriture de l'image vers {output_path}")
        
        with open(output_path, 'wb') as f:
            f.write(blob)
            
        logger.info(f"Image extraite avec succès : {output_path}")
        return output_path
    except Exception as e:
        logger.error(f"Erreur lors de l'extraction de l'image {filename}: {str(e)}")
        return None

def get_used_layouts_from_rels(pptx_path):
    """Récupère la liste des layouts utilisés en analysant les fichiers de relations."""
    used_layouts = set()
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            # Parcours des fichiers de relations des slides
            for file_info in zip_ref.infolist():
                if file_info.filename.startswith('ppt/slides/_rels/slide') and file_info.filename.endswith('.xml.rels'):
                    # Lecture du fichier de relations
                    rels_content = zip_ref.read(file_info.filename)
                    root = ET.fromstring(rels_content)
                    
                    # Recherche des références aux layouts
                    for rel in root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                        target = rel.get('Target', '')
                        if 'slideLayouts/slideLayout' in target:
                            # Extraction de l'ID du layout : Target="../slideLayouts/slideLayout88.xml"
                            layout_id = target.split('slideLayouts/slideLayout')[1].split('.')[0]
                            used_layouts.add(layout_id)
                            logger.debug(f"Layout utilisé détecté dans {file_info.filename}: {target}")
        
        logger.info(f"Layouts utilisés détectés : {sorted(used_layouts)}")
        return used_layouts
    except Exception as e:
        logger.error(f"Erreur lors de la lecture des relations : {str(e)}")
        return set()

def get_layout_images(pptx_path, used_layout_ids):
    """Récupère les images utilisées dans les layouts en analysant les fichiers de relations."""
    layout_images = {}
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            # Parcours des fichiers de relations des layouts
            for layout_id in used_layout_ids:
                rels_path = f'ppt/slideLayouts/_rels/slideLayout{layout_id}.xml.rels'
                if rels_path in zip_ref.namelist():
                    # Lecture du fichier de relations
                    rels_content = zip_ref.read(rels_path)
                    root = ET.fromstring(rels_content)
                    
                    # Recherche des références aux images
                    for rel in root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                        target = rel.get('Target', '')
                        if target.startswith('../media/'):
                            image_name = os.path.basename(target)
                            if layout_id not in layout_images:
                                layout_images[layout_id] = set()
                            layout_images[layout_id].add(image_name)
                            logger.debug(f"Image {image_name} trouvée dans le layout {layout_id} via {rels_path}")
        
        return layout_images
    except Exception as e:
        logger.error(f"Erreur lors de la récupération des images des layouts : {str(e)}")
        return {}

def is_slide_hidden(pptx_path, slide_index):
    """Vérifie si une diapositive est masquée en lisant son fichier XML."""
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            # Lecture du fichier XML de la diapositive
            slide_xml_path = f'ppt/slides/slide{slide_index}.xml'
            if slide_xml_path in zip_ref.namelist():
                slide_xml = zip_ref.read(slide_xml_path)
                root = ET.fromstring(slide_xml)
                # Vérification de l'attribut show
                show_attr = root.get('show')
                return show_attr == '0'
    except Exception as e:
        logger.error(f"Erreur lors de la lecture du fichier XML de la diapositive {slide_index}: {str(e)}")
    return False

def analyze_pptx(file_path):
    logger.info(f"Début de l'analyse du fichier : {file_path}")
    prs = Presentation(file_path)
    
    # Récupération des layouts utilisés depuis les fichiers de relations
    used_layout_ids = get_used_layouts_from_rels(file_path)
    logger.info(f"Layouts utilisés détectés : {sorted(used_layout_ids)}")
    
    # Récupération des images des layouts utilisés
    layout_images = get_layout_images(file_path, used_layout_ids)
    
    # Création d'un ensemble des noms de fichiers des images utilisées dans les layouts
    layout_used_filenames = set()
    for layout_id, images in layout_images.items():
        layout_used_filenames.update(images)
        logger.info(f"Images du layout {layout_id} : {sorted(images)}")
    
    logger.info(f"Total des images utilisées dans les layouts : {len(layout_used_filenames)}")
    
    # Analyse des diapositives
    slide_images = {}
    for slide_index, slide in enumerate(prs.slides):
        slide_shapes = get_images_from_shapes(slide.shapes, file_path, slide_index + 1)
        for blob, info in slide_shapes.items():
            if blob not in slide_images:
                slide_images[blob] = info
            slide_images[blob]["used"] = True
            # Vérification de la propriété show dans le XML
            slide_images[blob]["is_hidden"] = is_slide_hidden(file_path, slide_index + 1)
    
    # Analyse des masters et layouts pour obtenir toutes les images
    all_images = {}
    for slide_master in prs.slide_masters:
        for layout in slide_master.slide_layouts:
            layout_id = str(layout.slide_id) if hasattr(layout, "slide_id") else str(layout)
            layout_shapes = get_images_from_shapes(layout.shapes, file_path)
            all_images.update(layout_shapes)
        
        master_shapes = get_images_from_shapes(slide_master.shapes, file_path)
        all_images.update(master_shapes)
    
    # Séparation des images utilisées et non utilisées
    unused_images = {
        blob: info for blob, info in all_images.items() 
        if not info.get("used", False) and info['filename'] not in layout_used_filenames
    }
    
    # Images rognées (uniquement celles des slides)
    used_and_cropped_images = {blob: info for blob, info in slide_images.items() if info["cropped"]}
    
    # Calcul des statistiques
    cropped_count = len(used_and_cropped_images)
    total_used_images = len(slide_images)
    percentage_cropped = (cropped_count / total_used_images * 100) if total_used_images > 0 else 0
    
    # Calcul du poids total des parties rognées
    total_cropped_size = sum(len(blob) * info["crop_ratio"] for blob, info in used_and_cropped_images.items())
    
    # Calcul du poids total des images inutilisées
    total_unused_size = sum(len(blob) for blob in unused_images.keys())

    # Analyse des dimensions physiques des images utilisées
    logger.debug("\nAnalyse des images utilisées :")
    used_images_info = []
    for blob, info in slide_images.items():
        try:
            # Conversion du blob en image PIL
            image = Image.open(io.BytesIO(blob))
            width, height = image.size
            
            # Ajout des informations dans la liste
            used_images_info.append({
                'filename': info['filename'],
                'size_kb': len(blob) / 1024,
                'slide_index': info['slide_index'],
                'dimensions': f"{width}x{height}",
                'is_cropped': info.get('cropped', False),
                'is_hidden': info.get('is_hidden', False),
                'crop_info': f"L={info.get('crop_left', 0)*100:.1f}%, T={info.get('crop_top', 0)*100:.1f}%, R={info.get('crop_right', 0)*100:.1f}%, B={info.get('crop_bottom', 0)*100:.1f}%" if info.get('cropped', False) else "Non rognée"
            })
            
        except Exception as e:
            logger.error(f"Erreur lors de l'analyse des dimensions de l'image {info['filename']}: {str(e)}")
    
    # Récupération des informations détaillées des layouts
    layout_info = get_layout_info(file_path, used_layout_ids)

    logger.info(f"Analyse terminée - Images non utilisées : {len(unused_images)}, Images rognées : {percentage_cropped:.2f}%")
    return file_path, unused_images, percentage_cropped, total_cropped_size, total_unused_size, used_and_cropped_images, used_layout_ids, layout_images, used_images_info, layout_info, layout_used_filenames

def get_media_files_from_pptx(pptx_path):
    """Extrait la liste des fichiers média et leurs relations depuis le PPTX."""
    media_files = {}
    with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
        # Lecture du fichier [Content_Types].xml pour trouver les relations
        content_types = zip_ref.read('[Content_Types].xml')
        root = ET.fromstring(content_types)
        
        # Lecture du fichier _rels/.rels pour trouver les relations principales
        rels = zip_ref.read('_rels/.rels')
        rels_root = ET.fromstring(rels)
        
        # Lecture du fichier ppt/_rels/presentation.xml.rels pour trouver les relations des médias
        ppt_rels = zip_ref.read('ppt/_rels/presentation.xml.rels')
        ppt_rels_root = ET.fromstring(ppt_rels)
        
        # Parcours des fichiers dans ppt/media
        for file_info in zip_ref.infolist():
            if file_info.filename.startswith('ppt/media/'):
                media_files[file_info.filename] = {
                    'size': file_info.file_size,
                    'used': False
                }
    
    return media_files

def create_light_version(file_path, unused_images):
    """Crée une version allégée du fichier PowerPoint en supprimant les images inutilisées."""
    try:
        logger.info(f"Création de la version allégée pour {file_path}")
        
        # Création du nom du fichier de sortie
        base_name = os.path.splitext(file_path)[0]
        output_path = f"{base_name}_light.pptx"
        
        # Création d'un répertoire temporaire
        with tempfile.TemporaryDirectory() as temp_dir:
            logger.debug(f"Création du répertoire temporaire : {temp_dir}")
            
            # Extraction du PPTX original
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
            
            # Analyse des fichiers média
            media_files = get_media_files_from_pptx(file_path)
            
            # Identification des fichiers média à supprimer
            files_to_remove = []
            for blob, info in unused_images.items():
                # Recherche du fichier média correspondant
                for media_file in media_files:
                    if media_file.endswith(f"/{info['filename']}"):
                        files_to_remove.append(media_file)
                        logger.debug(f"Fichier média à supprimer : {media_file}")
            
            # Suppression des fichiers média inutilisés
            for file_to_remove in files_to_remove:
                file_path = os.path.join(temp_dir, file_to_remove)
                if os.path.exists(file_path):
                    os.remove(file_path)
                    logger.debug(f"Suppression du fichier : {file_to_remove}")
            
            # Création du nouveau PPTX
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                for root, dirs, files in os.walk(temp_dir):
                    for file in files:
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, temp_dir)
                        zip_out.write(file_path, arcname)
                        logger.debug(f"Ajout du fichier dans l'archive : {arcname}")
        
        logger.info(f"Version allégée créée avec succès : {output_path}")
        return output_path
    except Exception as e:
        logger.error(f"Erreur lors de la création de la version allégée : {str(e)}")
        return None

def on_create_light_version():
    """Gestionnaire d'événement pour le bouton de création de version allégée"""
    if not hasattr(on_create_light_version, 'last_file_path') or not on_create_light_version.last_file_path:
        result_text.delete(1.0, "end")
        result_text.insert("end", "Veuillez d'abord analyser un fichier PowerPoint.")
        return
        
    try:
        output_path = create_light_version(on_create_light_version.last_file_path, on_create_light_version.last_unused_images)
        if output_path:
            result_text.delete(1.0, "end")
            result_text.insert("end", f"Version allégée créée avec succès :\n{output_path}")
        else:
            result_text.delete(1.0, "end")
            result_text.insert("end", "Erreur lors de la création de la version allégée.")
    except Exception as e:
        logger.error(f"Erreur lors de la création de la version allégée : {str(e)}")
        result_text.delete(1.0, "end")
        result_text.insert("end", f"Erreur : {str(e)}")

def get_layout_info(pptx_path, used_layout_ids):
    """Récupère les informations détaillées des layouts utilisés."""
    layout_info = {}
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            # Lecture du fichier presentation.xml pour obtenir les noms des layouts
            presentation_xml = zip_ref.read('ppt/presentation.xml')
            root = ET.fromstring(presentation_xml)
            
            # Parcours des layouts dans presentation.xml
            for layout in root.findall('.//p:sldLayout', {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}):
                layout_id = layout.get('id', '')
                if layout_id in used_layout_ids:
                    name = layout.get('name', f'Layout {layout_id}')
                    type_attr = layout.get('type', 'unknown')
                    layout_info[layout_id] = {
                        'name': name,
                        'type': type_attr
                    }
        
        return layout_info
    except Exception as e:
        logger.error(f"Erreur lors de la récupération des informations des layouts : {str(e)}")
        return {}

def crop_image(blob, info):
    """Convertit l'image à 150 DPI si sa résolution est supérieure."""
    try:
        # Conversion du blob en image PIL
        image = Image.open(io.BytesIO(blob))
        
        # Récupération de la résolution actuelle
        current_dpi = image.info.get('dpi', (72, 72))[0]  # Par défaut 72 DPI si non spécifié
        
        # Si la résolution est supérieure à 150 DPI, on convertit
        if current_dpi > 150:
            logger.debug(f"Conversion de l'image {info['filename']} de {current_dpi} DPI à 150 DPI")
            
            # Calcul des nouvelles dimensions pour maintenir la taille physique
            width, height = image.size
            new_width = int(width * (150 / current_dpi))
            new_height = int(height * (150 / current_dpi))
            
            # Redimensionnement de l'image
            resized_image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            # Conversion en bytes avec une résolution de 150 DPI
            output = io.BytesIO()
            resized_image.save(output, format=image.format, quality=95, dpi=(150, 150))
            return output.getvalue()
        else:
            logger.debug(f"Image {info['filename']} déjà en résolution inférieure ou égale à 150 DPI ({current_dpi} DPI)")
            return blob
            
    except Exception as e:
        logger.error(f"Erreur lors de la conversion de l'image {info.get('filename', 'inconnue')}: {str(e)}")
        return blob

def emu_to_percentage(emu_value):
    """Convertit une valeur EMU en pourcentage (0-1)."""
    # 100000 EMU = 1%
    return float(emu_value) / 10000000

def remove_crop_info_from_slide_xml(temp_dir, slide_index, cropped_images):
    """Supprime les informations de rognage et ajuste les dimensions des formes dans le fichier XML de la diapositive."""
    try:
        # Chemin vers le fichier XML de la diapositive
        slide_xml_path = os.path.join(temp_dir, 'ppt', 'slides', f'slide{slide_index}.xml')
        if not os.path.exists(slide_xml_path):
            logger.debug(f"Fichier XML de la diapositive {slide_index} non trouvé")
            return

        # Lecture du fichier XML
        tree = ET.parse(slide_xml_path)
        root = tree.getroot()

        # Espace de noms pour les éléments de forme
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
              'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        # Recherche de tous les éléments de forme avec des images
        for pic in root.findall('.//p:pic', ns):
            # Recherche des attributs de rognage dans blipFill
            blipFill = pic.find('.//a:blipFill', ns)
            if blipFill is not None:
                # Récupération de l'ID de l'image
                embed = blipFill.find('.//a:blip', ns)
                if embed is not None:
                    rId = embed.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if rId:
                        # Recherche du nom de fichier correspondant
                        image_name = None
                        for blob, info in cropped_images.items():
                            if info.get('slide_index') == slide_index:
                                image_name = info['filename']
                                break
                        
                        if image_name:
                            # Recherche de l'élément srcRect pour les valeurs EMU
                            srcRect = blipFill.find('.//a:srcRect', ns)
                            if srcRect is not None:
                                # Conversion des valeurs EMU en pourcentages
                                crop_values = {}
                                for attr in ['l', 't', 'r', 'b']:
                                    emu_value = srcRect.get(attr)
                                    if emu_value is not None:
                                        crop_values[attr] = emu_to_percentage(emu_value)
                                        logger.debug(f"Valeur EMU convertie pour {attr}: {emu_value} -> {crop_values[attr]:.4f}")
                                
                                # Mise à jour des informations de rognage
                                if crop_values:
                                    info['crop_left'] = crop_values.get('l', 0)
                                    info['crop_top'] = crop_values.get('t', 0)
                                    info['crop_right'] = crop_values.get('r', 0)
                                    info['crop_bottom'] = crop_values.get('b', 0)
                                    info['crop_ratio'] = sum(crop_values.values()) / len(crop_values)
                                    
                                    # Recherche de l'élément xfrm pour ajuster les dimensions
                                    xfrm = pic.find('.//a:xfrm', ns)
                                    if xfrm is not None:
                                        ext = xfrm.find('.//a:ext', ns)
                                        if ext is not None:
                                            current_width = int(ext.get('cx', 0))
                                            current_height = int(ext.get('cy', 0))
                                            
                                            # Calcul des nouvelles dimensions
                                            new_width = int(current_width * (1 - info['crop_left'] - info['crop_right']))
                                            new_height = int(current_height * (1 - info['crop_top'] - info['crop_bottom']))
                                            
                                            # Mise à jour des dimensions
                                            ext.set('cx', str(new_width))
                                            ext.set('cy', str(new_height))
                                            logger.debug(f"Ajustement des dimensions pour l'image {image_name} : "
                                                       f"{current_width}x{current_height} -> {new_width}x{new_height}")
                
                # Suppression des attributs de rognage dans blipFill
                for attr in ['cropLeft', 'cropTop', 'cropRight', 'cropBottom']:
                    if blipFill.get(attr) is not None:
                        logger.debug(f"Suppression de l'attribut {attr} dans blipFill de la diapositive {slide_index}")
                        blipFill.attrib.pop(attr)
                
                # Suppression des attributs de rognage dans srcRect
                srcRect = blipFill.find('.//a:srcRect', ns)
                if srcRect is not None:
                    # Suppression directe de tous les attributs de rognage
                    for attr in ['l', 't', 'r', 'b']:
                        if attr in srcRect.attrib:
                            logger.debug(f"Suppression de l'attribut {attr} dans srcRect de la diapositive {slide_index}")
                            srcRect.attrib.pop(attr)
                    
                    # Si srcRect n'a plus d'attributs, le retirer complètement
                    if not srcRect.attrib:
                        parent = srcRect.getparent()
                        if parent is not None:
                            parent.remove(srcRect)
                            logger.debug(f"Suppression de l'élément srcRect vide dans la diapositive {slide_index}")

        # Sauvegarde du fichier XML modifié
        tree.write(slide_xml_path, encoding='UTF-8', xml_declaration=True)
        logger.debug(f"Informations de rognage supprimées et dimensions ajustées pour la diapositive {slide_index}")

    except Exception as e:
        logger.error(f"Erreur lors de la modification de la diapositive {slide_index}: {str(e)}")

def update_pptx_with_cropped_images(file_path, cropped_images):
    """Met à jour le fichier PPTX avec les images rognées."""
    try:
        # Création du nom du fichier de sortie
        base_name = os.path.splitext(file_path)[0]
        output_path = f"{base_name}_cropped.pptx"
        
        # Création d'un répertoire temporaire
        with tempfile.TemporaryDirectory() as temp_dir:
            logger.debug(f"Création du répertoire temporaire : {temp_dir}")
            
            # Extraction du PPTX original
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
            
            # Mise à jour des images rognées
            processed_slides = set()
            for blob, info in cropped_images.items():
                if info["cropped"]:
                    # Recherche du fichier image dans le répertoire temporaire
                    image_path = os.path.join(temp_dir, 'ppt', 'media', info['filename'])
                    if os.path.exists(image_path):
                        # Rogne l'image avec les informations stockées
                        cropped_blob = crop_image(blob, info)
                        # Écrit l'image rognée
                        with open(image_path, 'wb') as f:
                            f.write(cropped_blob)
                        logger.debug(f"Image rognée mise à jour : {info['filename']}")
                        
                        # Suppression des informations de rognage dans le XML de la diapositive
                        if info.get('slide_index') is not None:
                            processed_slides.add(info['slide_index'])
            
            # Suppression des informations de rognage et ajustement des dimensions pour toutes les diapositives traitées
            for slide_index in processed_slides:
                remove_crop_info_from_slide_xml(temp_dir, slide_index, cropped_images)
            
            # Création du nouveau PPTX
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                for root, dirs, files in os.walk(temp_dir):
                    for file in files:
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, temp_dir)
                        zip_out.write(file_path, arcname)
                        logger.debug(f"Ajout du fichier dans l'archive : {arcname}")
        
        logger.info(f"Version avec images rognées créée : {output_path}")
        return output_path
    except Exception as e:
        logger.error(f"Erreur lors de la mise à jour du PPTX : {str(e)}")
        return None

def on_create_cropped_version():
    """Gestionnaire d'événement pour le bouton de création de version avec images rognées"""
    if not hasattr(on_create_cropped_version, 'last_file_path') or not on_create_cropped_version.last_file_path:
        result_text.delete(1.0, "end")
        result_text.insert("end", "Veuillez d'abord analyser un fichier PowerPoint.")
        return
        
    try:
        output_path = update_pptx_with_cropped_images(
            on_create_cropped_version.last_file_path,
            on_create_cropped_version.last_cropped_images
        )
        if output_path:
            result_text.delete(1.0, "end")
            result_text.insert("end", f"Version avec images rognées créée :\n{output_path}")
        else:
            result_text.delete(1.0, "end")
            result_text.insert("end", "Erreur lors de la création de la version avec images rognées.")
    except Exception as e:
        logger.error(f"Erreur lors de la création de la version avec images rognées : {str(e)}")
        result_text.delete(1.0, "end")
        result_text.insert("end", f"Erreur : {str(e)}")

def on_drop(event):
    # Nettoyage du chemin de fichier
    file_path = event.data.strip('{}')
    logger.info(f"Tentative d'analyse du fichier : {file_path}")
    
    if os.path.isfile(file_path) and file_path.lower().endswith(".pptx"):
        logger.info(f"Fichier valide détecté : {file_path}")
        try:
            file_name, unused_images, percentage_cropped, total_cropped_size, total_unused_size, cropped_images, used_layout_ids, layout_images, used_images_info, layout_info, layout_used_filenames = analyze_pptx(file_path)
            
            # Stockage des informations pour les boutons
            on_create_light_version.last_file_path = file_path
            on_create_light_version.last_unused_images = unused_images
            on_create_cropped_version.last_file_path = file_path
            on_create_cropped_version.last_cropped_images = cropped_images

            # Nettoyage de l'interface
            result_text.delete(1.0, "end")
            unused_tree.delete(*unused_tree.get_children())
            tree.delete(*tree.get_children())
            layout_tree.delete(*layout_tree.get_children())

            # Affichage des informations générales
            result_text.insert("end", f"Nom du fichier : {file_name}\n"
                                    f"Nombre total d'images non utilisées : {len(unused_images)}\n"
                                    f"Poids total des images non utilisées : {total_unused_size/1024:.2f} KB\n"
                                    f"Pourcentage d'images utilisées qui sont rognées : {percentage_cropped:.2f}%\n"
                                    f"Poids total estimé des parties rognées : {total_cropped_size/1024:.2f} KB\n\n")

            # Configuration des colonnes du tableau pour les images non utilisées
            unused_tree['columns'] = ('slide', 'visible', 'filename', 'size', 'dimensions', 'crop_info')
            unused_tree.column('#0', width=0, stretch=False)
            unused_tree.column('slide', width=50, anchor='center')
            unused_tree.column('visible', width=60, anchor='center')
            unused_tree.column('filename', width=200)
            unused_tree.column('size', width=80, anchor='e')
            unused_tree.column('dimensions', width=100, anchor='center')
            unused_tree.column('crop_info', width=200)

            # Configuration des en-têtes pour les images non utilisées
            unused_tree.heading('slide', text='Slide', command=lambda: sort_treeview(unused_tree, 'slide', False))
            unused_tree.heading('visible', text='Visible', command=lambda: sort_treeview(unused_tree, 'visible', False))
            unused_tree.heading('filename', text='Nom du fichier', command=lambda: sort_treeview(unused_tree, 'filename', False))
            unused_tree.heading('size', text='Taille (KB)', command=lambda: sort_treeview(unused_tree, 'size', False))
            unused_tree.heading('dimensions', text='Dimensions', command=lambda: sort_treeview(unused_tree, 'dimensions', False))
            unused_tree.heading('crop_info', text='Rognage', command=lambda: sort_treeview(unused_tree, 'crop_info', False))

            # Configuration des colonnes du tableau pour les images des layouts
            layout_tree['columns'] = ('layout', 'used', 'filename', 'size', 'dimensions', 'crop_info')
            layout_tree.column('#0', width=0, stretch=False)
            layout_tree.column('layout', width=100, anchor='center')
            layout_tree.column('used', width=60, anchor='center')
            layout_tree.column('filename', width=200)
            layout_tree.column('size', width=80, anchor='e')
            layout_tree.column('dimensions', width=100, anchor='center')
            layout_tree.column('crop_info', width=200)

            # Configuration des en-têtes pour les images des layouts
            layout_tree.heading('layout', text='Layout', command=lambda: sort_treeview(layout_tree, 'layout', False))
            layout_tree.heading('used', text='Utilisé', command=lambda: sort_treeview(layout_tree, 'used', False))
            layout_tree.heading('filename', text='Nom du fichier', command=lambda: sort_treeview(layout_tree, 'filename', False))
            layout_tree.heading('size', text='Taille (KB)', command=lambda: sort_treeview(layout_tree, 'size', False))
            layout_tree.heading('dimensions', text='Dimensions', command=lambda: sort_treeview(layout_tree, 'dimensions', False))
            layout_tree.heading('crop_info', text='Rognage', command=lambda: sort_treeview(layout_tree, 'crop_info', False))

            # Configuration des colonnes du tableau pour les images utilisées
            tree['columns'] = ('slide', 'visible', 'filename', 'size', 'dimensions', 'crop_info')
            tree.column('#0', width=0, stretch=False)
            tree.column('slide', width=50, anchor='center')
            tree.column('visible', width=60, anchor='center')
            tree.column('filename', width=200)
            tree.column('size', width=80, anchor='e')
            tree.column('dimensions', width=100, anchor='center')
            tree.column('crop_info', width=200)

            # Configuration des en-têtes pour les images utilisées
            tree.heading('slide', text='Slide', command=lambda: sort_treeview(tree, 'slide', False))
            tree.heading('visible', text='Visible', command=lambda: sort_treeview(tree, 'visible', False))
            tree.heading('filename', text='Nom du fichier', command=lambda: sort_treeview(tree, 'filename', False))
            tree.heading('size', text='Taille (KB)', command=lambda: sort_treeview(tree, 'size', False))
            tree.heading('dimensions', text='Dimensions', command=lambda: sort_treeview(tree, 'dimensions', False))
            tree.heading('crop_info', text='Rognage', command=lambda: sort_treeview(tree, 'crop_info', False))

            # Affichage des images non utilisées
            if unused_images:
                # Trier les images par taille décroissante
                sorted_images = sorted(unused_images.items(), key=lambda x: len(x[0]), reverse=True)
                for blob, info in sorted_images:
                    try:
                        image = Image.open(io.BytesIO(blob))
                        width, height = image.size
                        unused_tree.insert('', 'end', values=(
                            "N/A",
                            "N/A",
                            info['filename'],
                            f"{len(blob)/1024:.2f}",
                            f"{width}x{height}",
                            "Non rognée"
                        ))
                    except Exception as e:
                        logger.error(f"Erreur lors de l'analyse des dimensions de l'image {info['filename']}: {str(e)}")
            else:
                unused_tree.insert('', 'end', values=("N/A", "N/A", "Aucune image non utilisée trouvée", "N/A", "N/A", "N/A"))

            # Affichage des images des layouts
            layout_total_size = 0
            for layout_id, images in layout_images.items():
                try:
                    layout_num = int(layout_id)  # Conversion en entier
                    layout_name = f"Layout {layout_num:03d}"  # Format "Layout 999"
                except ValueError:
                    layout_name = f"Layout {layout_id}"  # Fallback si la conversion échoue
                
                # Vérification si le layout est utilisé dans les slides
                is_used = layout_id in used_layout_ids
                
                for image_name in sorted(images):
                    try:
                        with zipfile.ZipFile(file_path, 'r') as zip_ref:
                            image_path = f'ppt/media/{image_name}'
                            if image_path in zip_ref.namelist():
                                image_data = zip_ref.read(image_path)
                                image = Image.open(io.BytesIO(image_data))
                                width, height = image.size
                                size_kb = len(image_data) / 1024
                                layout_total_size += size_kb
                                layout_tree.insert('', 'end', values=(
                                    layout_name,
                                    "Oui" if is_used else "Non",
                                    image_name,
                                    f"{size_kb:.2f}",
                                    f"{width}x{height}",
                                    "Non rognée"
                                ))
                    except Exception as e:
                        logger.error(f"Erreur lors de l'analyse des dimensions de l'image {image_name}: {str(e)}")

            # Ajout des images dans le tableau des images utilisées
            for img_info in used_images_info:
                # Détermination si l'image est dans un layout ou une slide
                is_layout = img_info['filename'] in layout_used_filenames
                if not is_layout:  # On n'affiche que les images des slides ici
                    slide_info = f"{img_info['slide_index']:03d} S"
                    tree.insert('', 'end', values=(
                        slide_info,
                        "Non" if img_info.get('is_hidden', False) else "Oui",
                        img_info['filename'],
                        f"{img_info['size_kb']:.2f}",
                        img_info['dimensions'],
                        img_info['crop_info']
                    ))

            # Mise à jour des titres avec les poids totaux
            total_used_size = sum(img_info['size_kb'] for img_info in used_images_info if img_info['filename'] not in layout_used_filenames)
            unused_label.config(text=f"Images supprimables : poids {total_unused_size/1024:.2f} Ko")
            used_label.config(text=f"Images des slides : poids {total_used_size:.2f} Ko")
            layout_label.config(text=f"Images des layouts : poids {layout_total_size:.2f} Ko")

            # Affichage des layouts utilisés
            result_text.insert("end", f"\nLayouts utilisés ({len(used_layout_ids)}) :\n")
            for layout_id in sorted(used_layout_ids):
                info = layout_info.get(layout_id, {})
                name = info.get('name', f'Layout {layout_id}')
                type_attr = info.get('type', 'unknown')
                result_text.insert("end", f"\n- Layout {layout_id} : {name} (Type: {type_attr})")
                
                if layout_id in layout_images and layout_images[layout_id]:
                    result_text.insert("end", "\n  Images utilisées :")
                    for image_name in sorted(layout_images[layout_id]):
                        result_text.insert("end", f"\n    - {image_name}")

        except Exception as e:
            logger.error(f"Erreur lors de l'analyse du fichier : {str(e)}")
            result_text.delete(1.0, "end")
            result_text.insert("end", f"Erreur lors de l'analyse du fichier : {str(e)}")
    else:
        logger.warning(f"Fichier invalide déposé : {file_path}")
        result_text.delete(1.0, "end")
        result_text.insert("end", "Veuillez déposer un fichier .pptx valide.")

def sort_treeview(tree_widget, col, reverse):
    """Fonction pour trier le tableau selon la colonne sélectionnée."""
    # Récupération de toutes les valeurs
    l = [(tree_widget.set(k, col), k) for k in tree_widget.get_children('')]
    
    # Tri des valeurs
    if col == 'size':  # Colonne de taille en KB
        # Conversion des valeurs en nombres pour le tri
        l.sort(key=lambda x: float(x[0].split()[0]), reverse=reverse)
    else:
        l.sort(reverse=reverse)
    
    # Réorganisation des éléments
    for index, (val, k) in enumerate(l):
        tree_widget.move(k, '', index)
    
    # Inversion du tri pour le prochain clic
    tree_widget.heading(col, command=lambda: sort_treeview(tree_widget, col, not reverse))

def on_closing():
    """Fonction appelée lors de la fermeture de l'application"""
    logger.info("Fermeture de l'application")
    root.destroy()

def draw_dashed_rectangle(draw, xy, outline, width=1, dash_length=5, gap_length=5):
    """Dessine un rectangle en pointillés sur une image PIL."""
    x1, y1 = xy[0]
    x2, y2 = xy[1]
    # Haut
    for x in range(x1, x2, dash_length + gap_length):
        draw.line([(x, y1), (min(x + dash_length, x2), y1)], fill=outline, width=width)
    # Bas
    for x in range(x1, x2, dash_length + gap_length):
        draw.line([(x, y2), (min(x + dash_length, x2), y2)], fill=outline, width=width)
    # Gauche
    for y in range(y1, y2, dash_length + gap_length):
        draw.line([(x1, y), (x1, min(y + dash_length, y2))], fill=outline, width=width)
    # Droite
    for y in range(y1, y2, dash_length + gap_length):
        draw.line([(x2, y), (x2, min(y + dash_length, y2))], fill=outline, width=width)

def show_image_with_frame(filename, window_title):
    """Affiche une image avec ses cadres dans une nouvelle fenêtre, limitée à 80% de la taille de l'écran, avec un cadre rouge clignotant plus épais."""
    try:
        image_window = Toplevel(root)
        image_window.title(window_title)
        screen_width = root.winfo_screenwidth()
        screen_height = root.winfo_screenheight()
        max_width = int(screen_width * 0.8)
        max_height = int(screen_height * 0.8)
        with zipfile.ZipFile(on_create_light_version.last_file_path, 'r') as zip_ref:
            image_path = f'ppt/media/{filename}'
            if image_path in zip_ref.namelist():
                image_data = zip_ref.read(image_path)
                image = Image.open(io.BytesIO(image_data))
                width, height = image.size
                scale = min(max_width / width, max_height / height, 1.0)
                if scale < 1.0:
                    new_width = int(width * scale)
                    new_height = int(height * scale)
                    image = image.resize((new_width, new_height), Image.LANCZOS)
                else:
                    new_width, new_height = width, height
                # Recherche des informations de rognage dans les données de l'image
                is_cropped = False
                crop_info = None
                for blob, info in on_create_cropped_version.last_cropped_images.items():
                    if info['filename'] == filename:
                        is_cropped = True
                        crop_info = info
                        break
                def update_frame_color(color1, color2, delay=400):
                    image_with_frame = image.copy()
                    draw = ImageDraw.Draw(image_with_frame)
                    if is_cropped and crop_info:
                        crop_left = int(new_width * crop_info['crop_left'])
                        crop_top = int(new_height * crop_info['crop_top'])
                        crop_right = int(new_width * (1 - crop_info['crop_right']))
                        crop_bottom = int(new_height * (1 - crop_info['crop_bottom']))
                        draw.rectangle([(crop_left, crop_top), (crop_right, crop_bottom)], outline=color1, width=4)
                        draw_dashed_rectangle(draw, [(0, 0), (new_width-1, new_height-1)], outline='blue', width=2)
                    else:
                        draw.rectangle([(0, 0), (new_width-1, new_height-1)], outline=color1, width=4)
                    photo = ImageTk.PhotoImage(image_with_frame)
                    label.config(image=photo)
                    label.image = photo
                    image_window.after(delay, lambda: update_frame_color(color2, color1, delay))
                photo = ImageTk.PhotoImage(image)
                label = Label(image_window, image=photo)
                label.image = photo
                label.pack()
                image_window.geometry(f"{new_width}x{new_height}")
                close_button = Button(image_window, text="Fermer", command=image_window.destroy)
                close_button.pack(pady=5)
                # Lancer le clignotement du cadre rouge (rouge <-> jaune)
                update_frame_color('red', 'yellow')
            else:
                Label(image_window, text="Image non trouvée").pack()
    except Exception as e:
        logger.error(f"Erreur lors de l'affichage de l'image : {str(e)}")
        messagebox.showerror("Erreur", f"Impossible d'afficher l'image : {str(e)}")

def show_image(event):
    """Affiche l'image sélectionnée dans une nouvelle fenêtre."""
    try:
        item = tree.selection()[0]
        filename = tree.item(item)['values'][2]
        show_image_with_frame(filename, f"Image : {filename}")
    except Exception as e:
        logger.error(f"Erreur lors de l'affichage de l'image : {str(e)}")
        messagebox.showerror("Erreur", f"Impossible d'afficher l'image : {str(e)}")

def show_unused_image(event):
    """Affiche l'image non utilisée sélectionnée dans une nouvelle fenêtre."""
    try:
        item = unused_tree.selection()[0]
        filename = unused_tree.item(item)['values'][2]
        show_image_with_frame(filename, f"Image non utilisée : {filename}")
    except Exception as e:
        logger.error(f"Erreur lors de l'affichage de l'image non utilisée : {str(e)}")
        messagebox.showerror("Erreur", f"Impossible d'afficher l'image : {str(e)}")

def show_layout_image(event):
    """Affiche l'image du layout sélectionnée dans une nouvelle fenêtre."""
    try:
        item = layout_tree.selection()[0]
        filename = layout_tree.item(item)['values'][2]
        show_image_with_frame(filename, f"Image du layout : {filename}")
    except Exception as e:
        logger.error(f"Erreur lors de l'affichage de l'image du layout : {str(e)}")
        messagebox.showerror("Erreur", f"Impossible d'afficher l'image : {str(e)}")

# Configuration de l'application
root = TkinterDnD.Tk()
root.title("Analyse PPTX")

# Message d'accueil
welcome_text = "Déposez un fichier PowerPoint (.pptx) ici pour l'analyser\n\n" + \
               "L'application va :\n" + \
               "- Identifier les images non utilisées\n" + \
               "- Calculer le pourcentage d'images rognées\n" + \
               "- Estimer le poids des parties rognées\n" + \
               "- Trier les images par taille décroissante\n" + \
               "- Permettre de créer une version allégée\n" + \
               "- Permettre de créer une version avec images rognées"

# Création du widget Text avec scrollbar pour les informations générales
text_frame = Text(root, wrap="word", width=60, height=10)
text_scrollbar = Scrollbar(text_frame, command=text_frame.yview)
text_frame.configure(yscrollcommand=text_scrollbar.set)

# Création du widget Treeview pour les images non utilisées
unused_frame = Frame(root)
unused_label = Label(unused_frame, text="Images supprimables : poids 0 Ko")
unused_tree = ttk.Treeview(unused_frame)
unused_scrollbar = Scrollbar(unused_frame, orient="vertical", command=unused_tree.yview)
unused_tree.configure(yscrollcommand=unused_scrollbar.set)

# Création du widget Treeview pour les images des layouts
layout_frame = Frame(root)
layout_label = Label(layout_frame, text="Images des layouts : poids 0 Ko")
layout_tree = ttk.Treeview(layout_frame)
layout_scrollbar = Scrollbar(layout_frame, orient="vertical", command=layout_tree.yview)
layout_tree.configure(yscrollcommand=layout_scrollbar.set)

# Création du widget Treeview pour les images utilisées
tree_frame = Frame(root)
used_label = Label(tree_frame, text="Images affichées : poids 0 Ko")
tree = ttk.Treeview(tree_frame)
tree_scrollbar = Scrollbar(tree_frame, orient="vertical", command=tree.yview)
tree.configure(yscrollcommand=tree_scrollbar.set)

# Création des boutons
button_frame = Frame(root)
close_button = Button(button_frame, text="Fermer", command=on_closing)
light_version_button = Button(button_frame, text="Créer version allégée", command=on_create_light_version)
cropped_version_button = Button(button_frame, text="Créer version rognée", command=on_create_cropped_version)
close_button.pack(side="top", padx=5, pady=5, fill="x")
cropped_version_button.pack(side="top", padx=5, pady=5, fill="x")
light_version_button.pack(side="top", padx=5, pady=5, fill="x")

# Placement des widgets principaux dans la fenêtre
main_content = Frame(root)
main_content.pack(side="left", fill="both", expand=True)
button_frame.pack(side="right", fill="y")

# Placement des widgets dans main_content
text_scrollbar.pack(side="right", fill="y")
text_frame.pack(side="top", fill="x", padx=5, pady=5)

unused_label.pack(side="top", anchor="w", padx=5)
unused_scrollbar.pack(side="right", fill="y")
unused_tree.pack(side="left", fill="both", expand=True)
unused_frame.pack(side="top", fill="x", padx=5, pady=5)

layout_label.pack(side="top", anchor="w", padx=5)
layout_scrollbar.pack(side="right", fill="y")
layout_tree.pack(side="left", fill="both", expand=True)
layout_frame.pack(side="top", fill="x", padx=5, pady=5)

used_label.pack(side="top", anchor="w", padx=5)
tree_scrollbar.pack(side="right", fill="y")
tree.pack(side="left", fill="both", expand=True)
tree_frame.pack(side="top", fill="both", expand=True, padx=5, pady=5)

# Configuration du texte
result_text = text_frame
result_text.insert("end", welcome_text)

# Configuration du glisser-déposer
root.drop_target_register(DND_FILES)
root.dnd_bind('<<Drop>>', on_drop)

# Configuration de la fermeture propre
root.protocol("WM_DELETE_WINDOW", on_closing)

# Configuration des événements de double-clic
tree.bind('<Double-1>', show_image)
unused_tree.bind('<Double-1>', show_unused_image)
layout_tree.bind('<Double-1>', show_layout_image)

logger.info("Démarrage de l'application")
root.mainloop()