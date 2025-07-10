import os
import tkinter as tk
from tkinter import ttk
from tkinterdnd2 import DND_FILES, TkinterDnD
from PIL import Image, ImageTk
import io
import PyPDF2
from pptx import Presentation
import logging
from logging_config import logger

class FileAnalyzer:
    def __init__(self):
        self.root = TkinterDnD.Tk()
        self.root.title("Analyseur de fichiers PDF et PPTX")
        self.root.geometry("800x600")
        
        # Configuration de la zone de dépôt
        self.drop_frame = ttk.LabelFrame(self.root, text="Déposez votre fichier ici")
        self.drop_frame.pack(fill="both", expand=True, padx=10, pady=5)
        
        self.drop_label = ttk.Label(self.drop_frame, text="Glissez-déposez un fichier PDF ou PPTX ici")
        self.drop_label.pack(pady=20)
        
        # Configuration de la zone de liste
        self.list_frame = ttk.LabelFrame(self.root, text="Images trouvées")
        self.list_frame.pack(fill="both", expand=True, padx=10, pady=5)
        
        # Création du Treeview
        self.tree = ttk.Treeview(self.list_frame, columns=("Type", "Page/Diapo", "Taille"), show="headings")
        self.tree.heading("Type", text="Type")
        self.tree.heading("Page/Diapo", text="Page/Diapo")
        self.tree.heading("Taille", text="Taille")
        
        # Ajout de la barre de défilement
        scrollbar = ttk.Scrollbar(self.list_frame, orient="vertical", command=self.tree.yview)
        self.tree.configure(yscrollcommand=scrollbar.set)
        
        # Placement des widgets
        self.tree.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # Configuration du glisser-déposer
        self.drop_label.drop_target_register(DND_FILES)
        self.drop_label.dnd_bind('<<Drop>>', self.on_drop)
        
        # Configuration du clic sur les éléments
        self.tree.bind('<Double-1>', self.show_image)
        
        # Stockage des images
        self.images = []
        
    def analyze_pdf(self, file_path):
        """Analyse un fichier PDF et extrait les images."""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    
                    if '/Resources' in page and '/XObject' in page['/Resources']:
                        x_objects = page['/Resources']['/XObject']
                        
                        for obj in x_objects:
                            if x_objects[obj]['/Subtype'] == '/Image':
                                try:
                                    image_data = x_objects[obj].get_data()
                                    image = Image.open(io.BytesIO(image_data))
                                    size = f"{image.size[0]}x{image.size[1]}"
                                    self.images.append({
                                        'type': 'PDF',
                                        'page': page_num + 1,
                                        'size': size,
                                        'data': image_data
                                    })
                                    self.tree.insert("", "end", values=("PDF", f"Page {page_num + 1}", size))
                                except Exception as e:
                                    logger.error(f"Erreur lors de l'extraction de l'image: {e}")
        except Exception as e:
            logger.error(f"Erreur lors de la lecture du PDF: {e}")
    
    def analyze_pptx(self, file_path):
        """Analyse un fichier PPTX et extrait les images."""
        try:
            prs = Presentation(file_path)
            
            for slide_index, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            image_data = shape.image.blob
                            image = Image.open(io.BytesIO(image_data))
                            size = f"{image.size[0]}x{image.size[1]}"
                            self.images.append({
                                'type': 'PPTX',
                                'page': slide_index + 1,
                                'size': size,
                                'data': image_data
                            })
                            self.tree.insert("", "end", values=("PPTX", f"Diapo {slide_index + 1}", size))
                        except Exception as e:
                            logger.error(f"Erreur lors de l'extraction de l'image: {e}")
        except Exception as e:
            logger.error(f"Erreur lors de la lecture du PPTX: {e}")
    
    def on_drop(self, event):
        """Gère l'événement de dépôt de fichier."""
        # Nettoyage du chemin de fichier
        file_path = event.data.strip('{}')
        
        # Vérification de l'extension du fichier
        _, ext = os.path.splitext(file_path.lower())
        
        # Nettoyage de la liste précédente
        self.tree.delete(*self.tree.get_children())
        self.images.clear()
        
        if ext == '.pdf':
            self.analyze_pdf(file_path)
        elif ext == '.pptx':
            self.analyze_pptx(file_path)
        else:
            logger.error(f"Format de fichier non supporté: {ext}")
    
    def show_image(self, event):
        """Affiche l'image sélectionnée dans une nouvelle fenêtre."""
        selection = self.tree.selection()
        if not selection:
            return
            
        item = self.tree.item(selection[0])
        index = self.tree.index(selection[0])
        
        if 0 <= index < len(self.images):
            image_data = self.images[index]['data']
            image = Image.open(io.BytesIO(image_data))
            
            # Création d'une nouvelle fenêtre
            window = tk.Toplevel(self.root)
            window.title(f"Image - {item['values'][0]} - {item['values'][1]}")
            
            # Redimensionnement de l'image pour l'affichage
            max_size = (800, 600)
            image.thumbnail(max_size, Image.Resampling.LANCZOS)
            
            # Conversion pour Tkinter
            photo = ImageTk.PhotoImage(image)
            
            # Affichage de l'image
            label = ttk.Label(window, image=photo)
            label.image = photo  # Garde une référence
            label.pack(padx=10, pady=10)
    
    def run(self):
        """Lance l'application."""
        self.root.mainloop()

if __name__ == "__main__":
    app = FileAnalyzer()
    app.run() 